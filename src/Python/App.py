from flask import Flask, request, send_file
from flask_cors import CORS
from pptx import Presentation
import io
import os

app = Flask(__name__)

# UBAH DISINI: Izinkan semua domain (termasuk GitHub Pages kamu) agar tidak kena blokir CORS
CORS(app)

# Tambahkan rute utama agar Render tidak 404 saat "bangun tidur"
@app.route('/')
def home():
    return "Server Asisten Gereja Aktif! Siap memproses PPT."

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.json  # Mengambil data slides dari React
    template_file = 'Template PowerPoint.pptx'
    
    # Dapatkan path absolut agar tidak error "file not found" di server
    base_path = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(base_path, template_file)
    
    try:
        # Memastikan file template ada
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Kesalahan Template: {e}")
        return {"error": f"Template tidak ditemukan di: {template_path}"}, 500

    # Melakukan looping untuk setiap slide yang dikirim dari form web
    for item in data['slides']:
        try:
            # Mengambil index layout dan validasi
            layout_idx = int(item['layout_idx'])
            if layout_idx >= len(prs.slide_layouts):
                continue
                
            layout_dipilih = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout_dipilih)
            
            # Mengurutkan placeholder agar pengisian judul dan isi tidak tertukar
            shapes = sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx)

            # Mengisi Judul (Placeholder pertama)
            if len(shapes) > 0 and item.get('judul'):
                shapes[0].text = item['judul']
            
            # Mengisi Isi/Konten (Placeholder kedua)
            if len(shapes) > 1 and item.get('isi'):
                tf = shapes[1].text_frame
                # Menangani baris baru agar tidak berantakan
                tf.text = item['isi']
                
        except Exception as e:
            print(f"Gagal memproses slide: {e}")
            continue

    # Menggunakan BytesIO agar file disimpan di RAM
    target_stream = io.BytesIO()
    prs.save(target_stream)
    target_stream.seek(0)

    return send_file(
        target_stream,
        as_attachment=True,
        download_name="Ibadah_Minggu.pptx",
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

if __name__ == '__main__':
    # Konfigurasi Port Dinamis untuk Render
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)